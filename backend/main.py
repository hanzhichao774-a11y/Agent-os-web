import os
import re
import json
import time
import sqlite3
import asyncio
import importlib.util
import inspect
from pathlib import Path
from datetime import datetime, timezone
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse
from pydantic import BaseModel
from agno.agent import Agent
from agno.db.sqlite import SqliteDb
from agno.knowledge.knowledge import Knowledge
from agno.knowledge.embedder.fastembed import FastEmbedEmbedder
from agno.knowledge.chunking.recursive import RecursiveChunking
from agno.vectordb.lancedb import LanceDb
from agno.knowledge.reranker.base import Reranker as BaseReranker
from agno.knowledge.document.base import Document as AgnoDocument
from agno.tools.pandas import PandasTools
from agno.tools.duckdb import DuckDbTools
from agno.tools.csv_toolkit import CsvTools
from agno.tools.file import FileTools
from agno.tools.file_generation import FileGenerationTools
from agno.tools.python import PythonTools
from agno.tools.calculator import CalculatorTools
from agno.guardrails.prompt_injection import PromptInjectionGuardrail

load_dotenv()

BASE_DIR = Path(__file__).parent
SKILLS_DIR = BASE_DIR / "skills"
KNOWLEDGE_DOCS_DIR = BASE_DIR / "knowledge" / "docs"
WORKSPACE_DIR = BASE_DIR / "workspace"

SKILLS_DIR.mkdir(exist_ok=True)
KNOWLEDGE_DOCS_DIR.mkdir(parents=True, exist_ok=True)
(BASE_DIR / "data").mkdir(exist_ok=True)
WORKSPACE_DIR.mkdir(exist_ok=True)

# ═══════════════════════════════════════════════════════════════════════════════
# 项目数据库（SQLite）
# ═══════════════════════════════════════════════════════════════════════════════

_PROJECTS_DB = str(BASE_DIR / "data" / "projects.db")


def _get_projects_conn() -> sqlite3.Connection:
    conn = sqlite3.connect(_PROJECTS_DB)
    conn.row_factory = sqlite3.Row
    return conn


def _init_projects_db():
    conn = _get_projects_conn()
    conn.execute("""
        CREATE TABLE IF NOT EXISTS projects (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            description TEXT NOT NULL DEFAULT '',
            status TEXT NOT NULL DEFAULT 'active',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS tasks (
            id TEXT PRIMARY KEY,
            project_id TEXT NOT NULL,
            name TEXT NOT NULL,
            sort_order INTEGER NOT NULL DEFAULT 0,
            created_at TEXT NOT NULL,
            FOREIGN KEY (project_id) REFERENCES projects(id) ON DELETE CASCADE
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS settings (
            key TEXT PRIMARY KEY,
            value TEXT NOT NULL,
            updated_at TEXT NOT NULL
        )
    """)
    conn.commit()
    count = conn.execute("SELECT COUNT(*) FROM projects").fetchone()[0]
    if count == 0:
        now = datetime.now(timezone.utc).isoformat()
        seed = [
            ("p1", "Q3 财报分析", "季度财务数据汇总与可视化", "active", now, now),
            ("p2", "供应链优化", "物流路径与库存优化方案", "active", now, now),
            ("p3", "客户流失预警", "高价值客户流失风险评估", "idle", now, now),
        ]
        conn.executemany(
            "INSERT INTO projects (id, name, description, status, created_at, updated_at) VALUES (?,?,?,?,?,?)",
            seed,
        )
        conn.commit()
    conn.close()


_init_projects_db()

app = FastAPI(title="Agent OS API", version="2.0.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ═══════════════════════════════════════════════════════════════════════════════
# 模型工厂
# ═══════════════════════════════════════════════════════════════════════════════

_PROVIDER_API_KEY_MAP = {
    "kimi": "KIMI_API_KEY",
    "minimax": "MINIMAX_API_KEY",
    "openai": "OPENAI_API_KEY",
}

_PROVIDER_BASE_URL_MAP = {
    "kimi": "https://api.moonshot.cn/v1",
    "minimax": "https://api.minimaxi.com/v1",
}

_PROVIDER_DEFAULT_MODEL = {
    "kimi": "kimi-k2.6",
    "minimax": "MiniMax-M2.7",
    "openai": "gpt-4o-mini",
    "custom": "qwen-plus",
}


def _get_llm_config() -> dict:
    """从 DB 读取 LLM 配置，无则回退到 .env。"""
    conn = _get_projects_conn()
    rows = conn.execute(
        "SELECT key, value FROM settings WHERE key LIKE 'llm_%'"
    ).fetchall()
    conn.close()
    db_cfg = {r["key"]: r["value"] for r in rows}
    if db_cfg.get("llm_provider"):
        return {
            "provider": db_cfg["llm_provider"],
            "model_id": db_cfg.get("llm_model_id", ""),
            "api_key": db_cfg.get("llm_api_key", ""),
            "base_url": db_cfg.get("llm_base_url", ""),
        }
    provider = os.getenv("MODEL_PROVIDER", "openai").lower()
    env_key = _PROVIDER_API_KEY_MAP.get(provider)
    return {
        "provider": provider,
        "model_id": os.getenv("MODEL_ID", ""),
        "api_key": os.getenv(env_key) if env_key else os.getenv("CUSTOM_API_KEY", ""),
        "base_url": os.getenv("CUSTOM_BASE_URL", ""),
    }

_COMPAT_ROLE_MAP = {
    "system": "system",
    "user": "user",
    "assistant": "assistant",
    "tool": "tool",
    "model": "assistant",
}

_PROVIDER_EXTRA_KWARGS: dict[str, dict] = {
    "kimi": {
        "extra_body": {"thinking": {"type": "disabled"}},
        "role_map": _COMPAT_ROLE_MAP,
    },
    "minimax": {"role_map": _COMPAT_ROLE_MAP},
    "custom": {"role_map": _COMPAT_ROLE_MAP},
}


def create_model():
    cfg = _get_llm_config()
    provider = cfg["provider"]
    model_id = cfg["model_id"] or _PROVIDER_DEFAULT_MODEL.get(provider, "gpt-4o-mini")
    api_key = cfg["api_key"]
    base_url = cfg["base_url"] or _PROVIDER_BASE_URL_MAP.get(provider)

    from agno.models.openai import OpenAIChat
    kwargs = dict(id=model_id, timeout=120)
    if api_key:
        kwargs["api_key"] = api_key
    if base_url:
        kwargs["base_url"] = base_url
    kwargs.update(_PROVIDER_EXTRA_KWARGS.get(provider, {}))
    return OpenAIChat(**kwargs)


# ═══════════════════════════════════════════════════════════════════════════════
# 技能加载器
# ═══════════════════════════════════════════════════════════════════════════════

_skill_registry: dict[str, dict] = {}


def _load_skill_module(path: Path) -> dict | None:
    """从 .py 文件加载一个技能，返回 {meta, run_fn} 或 None。"""
    try:
        spec = importlib.util.spec_from_file_location(path.stem, path)
        mod = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(mod)
        meta = getattr(mod, "SKILL_META", None)
        run_fn = getattr(mod, "run", None)
        if not meta or not callable(run_fn):
            return None
        sig = inspect.signature(run_fn)
        params = []
        for name, p in sig.parameters.items():
            ptype = "string"
            if p.annotation in (float, int):
                ptype = "number"
            params.append({"name": name, "type": ptype, "default": str(p.default) if p.default is not inspect.Parameter.empty else None})
        return {
            "id": path.stem,
            "meta": meta,
            "run_fn": run_fn,
            "params": params,
            "file": str(path),
        }
    except Exception as e:
        print(f"[WARN] 加载技能 {path.name} 失败: {e}")
        return None


def scan_skills():
    """扫描 skills/ 目录，刷新技能注册表。"""
    _skill_registry.clear()
    for f in sorted(SKILLS_DIR.glob("*.py")):
        if f.name.startswith("_"):
            continue
        skill = _load_skill_module(f)
        if skill:
            _skill_registry[skill["id"]] = skill
            print(f"[SKILL] 已加载: {skill['meta']['name']} ({skill['id']})")


scan_skills()


# ═══════════════════════════════════════════════════════════════════════════════
# ═══════════════════════════════════════════════════════════════════════════════
# Embedding / Reranker 配置
# ═══════════════════════════════════════════════════════════════════════════════

def _get_embedding_config() -> dict:
    """从 DB 读取 Embedding 配置，无则回退到 .env。"""
    conn = _get_projects_conn()
    rows = conn.execute(
        "SELECT key, value FROM settings WHERE key LIKE 'embedding_%'"
    ).fetchall()
    conn.close()
    db_cfg = {r["key"]: r["value"] for r in rows}
    if db_cfg.get("embedding_model_id"):
        return {
            "mode": "api",
            "model_id": db_cfg["embedding_model_id"],
            "api_key": db_cfg.get("embedding_api_key", ""),
            "base_url": db_cfg.get("embedding_base_url", ""),
            "dimensions": int(db_cfg.get("embedding_dimensions", "1024")),
        }
    env_model = os.getenv("EMBEDDING_MODEL_ID", "")
    if env_model:
        return {
            "mode": "api",
            "model_id": env_model,
            "api_key": os.getenv("EMBEDDING_API_KEY", ""),
            "base_url": os.getenv("EMBEDDING_BASE_URL", ""),
            "dimensions": int(os.getenv("EMBEDDING_DIMENSIONS", "1024")),
        }
    return {"mode": "local", "model_id": "", "api_key": "", "base_url": "", "dimensions": 512}


def _get_reranker_config() -> dict:
    """从 DB 读取 Reranker 配置，无则回退到 .env。"""
    conn = _get_projects_conn()
    rows = conn.execute(
        "SELECT key, value FROM settings WHERE key LIKE 'reranker_%'"
    ).fetchall()
    conn.close()
    db_cfg = {r["key"]: r["value"] for r in rows}
    if db_cfg.get("reranker_model_id"):
        return {
            "enabled": db_cfg.get("reranker_enabled", "true").lower() == "true",
            "model_id": db_cfg["reranker_model_id"],
            "api_key": db_cfg.get("reranker_api_key", ""),
            "base_url": db_cfg.get("reranker_base_url", ""),
            "top_n": int(db_cfg.get("reranker_top_n", "5")),
        }
    env_model = os.getenv("RERANKER_MODEL_ID", "")
    if env_model:
        return {
            "enabled": True,
            "model_id": env_model,
            "api_key": os.getenv("RERANKER_API_KEY", ""),
            "base_url": os.getenv("RERANKER_BASE_URL", ""),
            "top_n": int(os.getenv("RERANKER_TOP_N", "5")),
        }
    return {"enabled": False, "model_id": "", "api_key": "", "base_url": "", "top_n": 5}


class OpenAICompatibleReranker(BaseReranker):
    """兼容 Jina/Cohere (/v1/rerank) 和 TEI (/rerank) 两种 rerank API 格式。"""
    model: str = "qw3-reranke-8b"
    api_key: str = ""
    base_url: str = ""
    top_n: int = 5

    def _call_rerank(self, query: str, documents: list[str]) -> list[dict]:
        import httpx
        headers = {"Content-Type": "application/json"}
        if self.api_key:
            headers["Authorization"] = f"Bearer {self.api_key}"

        base = self.base_url.rstrip("/")

        # 先尝试 Jina/Cohere 格式: POST /v1/rerank
        jina_url = f"{base}/v1/rerank"
        jina_body = {
            "model": self.model,
            "query": query,
            "documents": documents,
            "top_n": self.top_n,
        }
        try:
            resp = httpx.post(jina_url, json=jina_body, headers=headers, timeout=30)
            if resp.status_code != 404:
                resp.raise_for_status()
                data = resp.json()
                return data.get("results", [])
        except httpx.HTTPStatusError as e:
            if e.response.status_code != 404:
                raise
        except httpx.ConnectError:
            pass

        # Fallback: TEI 格式: POST /rerank
        tei_url = f"{base}/rerank"
        tei_body = {"query": query, "texts": documents, "truncate": True}
        resp = httpx.post(tei_url, json=tei_body, headers=headers, timeout=30)
        resp.raise_for_status()
        data = resp.json()
        if isinstance(data, list):
            return [{"index": item.get("index", i), "relevance_score": item.get("score", 0)} for i, item in enumerate(data)]
        return data.get("results", [])

    def rerank(self, query: str, documents: list[AgnoDocument]) -> list[AgnoDocument]:
        if not documents:
            return []
        try:
            doc_texts = [doc.content for doc in documents]
            results = self._call_rerank(query, doc_texts)
            scored: list[AgnoDocument] = []
            for r in results:
                idx = r.get("index", 0)
                if idx < len(documents):
                    doc = documents[idx]
                    doc.reranking_score = r.get("relevance_score", 0)
                    scored.append(doc)
            scored.sort(
                key=lambda x: x.reranking_score if x.reranking_score is not None else float("-inf"),
                reverse=True,
            )
            return scored[:self.top_n] if self.top_n else scored
        except Exception as e:
            print(f"[WARN] Reranker 调用失败，返回原始文档: {e}")
            return documents


# ═══════════════════════════════════════════════════════════════════════════════
# 知识库 (动态 Embedder + Reranker)
# ═══════════════════════════════════════════════════════════════════════════════

from agno.vectordb.search import SearchType

_knowledge: Knowledge | None = None
_vector_db: LanceDb | None = None


def _rebuild_knowledge():
    """根据 settings 配置动态重建 Knowledge（embedder + reranker）。"""
    global _knowledge, _vector_db

    emb_cfg = _get_embedding_config()
    rer_cfg = _get_reranker_config()

    if emb_cfg["mode"] == "api" and emb_cfg["model_id"]:
        from agno.knowledge.embedder.openai import OpenAIEmbedder
        emb_kwargs: dict = {"id": emb_cfg["model_id"]}
        if emb_cfg["dimensions"]:
            emb_kwargs["dimensions"] = emb_cfg["dimensions"]
        if emb_cfg["api_key"]:
            emb_kwargs["api_key"] = emb_cfg["api_key"]
        if emb_cfg["base_url"]:
            emb_kwargs["base_url"] = emb_cfg["base_url"]
        embedder = OpenAIEmbedder(**emb_kwargs)
        print(f"[KNOWLEDGE] 使用远程 Embedding: {emb_cfg['model_id']} @ {emb_cfg['base_url'] or 'default'}")
    else:
        embedder = FastEmbedEmbedder(id="BAAI/bge-small-zh-v1.5", dimensions=512)
        print("[KNOWLEDGE] 使用本地 FastEmbed: BAAI/bge-small-zh-v1.5")

    reranker = None
    if rer_cfg["enabled"] and rer_cfg["model_id"]:
        reranker = OpenAICompatibleReranker(
            model=rer_cfg["model_id"],
            api_key=rer_cfg["api_key"],
            base_url=rer_cfg["base_url"],
            top_n=rer_cfg["top_n"],
        )
        print(f"[KNOWLEDGE] 使用 Reranker: {rer_cfg['model_id']} @ {rer_cfg['base_url'] or 'default'}")
    else:
        print("[KNOWLEDGE] Reranker 未启用")

    _vector_db = LanceDb(
        uri=str(BASE_DIR / "data" / "lancedb"),
        table_name="knowledge",
        embedder=embedder,
        search_type=SearchType.hybrid,
        reranker=reranker,
    )

    _knowledge = Knowledge(
        vector_db=_vector_db,
        max_results=10,
    )
    print("[KNOWLEDGE] Knowledge 初始化完成")


_rebuild_knowledge()


def ingest_document(doc_name: str, text: str) -> int:
    """将文本通过 RecursiveChunking 分块后存入向量库。"""
    from agno.knowledge.reader.text_reader import TextReader
    reader = TextReader(
        chunk=True,
        chunk_size=1500,
        chunking_strategy=RecursiveChunking(chunk_size=1500, overlap=200),
    )
    _knowledge.insert(
        name=doc_name,
        text_content=text,
        reader=reader,
        upsert=True,
    )
    from agno.knowledge.document.base import Document
    doc = Document(content=text, name=doc_name)
    chunks = RecursiveChunking(chunk_size=1500, overlap=200).chunk(doc)
    return len(chunks)


def search_knowledge(query: str, top_k: int = 5) -> list[dict]:
    """向量语义检索。"""
    results = _knowledge.search(query, max_results=top_k)
    chunks = []
    for doc in results:
        meta = doc.meta_data or {}
        name = meta.get("content_name") or meta.get("name") or doc.name or "unknown"
        chunks.append({
            "doc_name": name,
            "chunk_text": doc.content,
        })
    return chunks


_uploaded_docs: dict[str, int] = {}


def _restore_uploaded_docs():
    """启动时扫描 KNOWLEDGE_DOCS_DIR，恢复文档名列表到内存。"""
    if not KNOWLEDGE_DOCS_DIR.exists():
        return
    for f in KNOWLEDGE_DOCS_DIR.iterdir():
        if f.is_file() and not f.name.startswith("."):
            _uploaded_docs.setdefault(f.name, -1)
    if _uploaded_docs:
        print(f"[KNOWLEDGE] 从磁盘恢复了 {len(_uploaded_docs)} 个文档记录: {', '.join(_uploaded_docs.keys())}")


_restore_uploaded_docs()


def list_documents() -> list[dict]:
    return [{"doc_name": k, "chunks": v} for k, v in _uploaded_docs.items()]


# ═══════════════════════════════════════════════════════════════════════════════
# Agno 内置工具工厂
# ═══════════════════════════════════════════════════════════════════════════════

def _make_builtin_tools(tool_names: list[str]) -> list:
    """根据工具名称列表实例化 Agno 内置 Toolkit。"""
    ws = str(WORKSPACE_DIR)
    factories = {
        "pandas": lambda: PandasTools(),
        "duckdb": lambda: DuckDbTools(db_path=str(WORKSPACE_DIR / "duckdb.db")),
        "csv": lambda: CsvTools(csvs=[ws]),
        "file": lambda: FileTools(base_dir=ws),
        "file_generation": lambda: FileGenerationTools(output_directory=ws),
        "python": lambda: PythonTools(base_dir=ws, restrict_to_base_dir=True),
        "calculator": lambda: CalculatorTools(),
    }
    tools = []
    for name in tool_names:
        factory = factories.get(name)
        if factory:
            tools.append(factory())
        else:
            print(f"[WARN] 未知内置工具: {name}")
    return tools


# ═══════════════════════════════════════════════════════════════════════════════
# Agent 配置
# ═══════════════════════════════════════════════════════════════════════════════

_agent_tools: dict[str, list[str]] = {}

AGENT_CONFIGS: dict[str, dict] = {
    "a1": {
        "name": "数据分析Agent",
        "avatar": "📊",
        "description": "擅长 SQL 查询、数据清洗、统计建模，可直接操作 CSV/Excel 数据",
        "capabilities": ["SQL", "Pandas", "DuckDB", "CSV", "可视化"],
        "builtin_tools": ["pandas", "duckdb", "csv", "file_generation"],
        "instructions": [
            "你是企业数据分析专家，擅长 SQL 查询、数据清洗和统计建模。",
            "你可以使用 Pandas 处理数据框、DuckDB 执行 SQL 查询、操作 CSV 文件。",
            "你还可以生成 CSV/JSON/PDF 格式的分析报告文件。",
            "回答要专业、简洁，需要时提供代码和分析结果。",
            "始终使用中文回答。",
        ],
    },
    "a2": {
        "name": "知识检索Agent",
        "avatar": "🔍",
        "description": "企业内部知识库问答与检索",
        "capabilities": ["RAG", "文档解析", "语义搜索"],
        "instructions": [
            "你是企业知识库检索专家。",
            "当用户提问时，系统会自动从知识库中检索相关文档片段并附在上下文中。",
            "你必须优先基于这些检索到的文档片段来回答问题。",
            "在回答中标注信息出处（如引用的文档名称）。",
            "如果知识库中没有找到相关内容，如实告知用户。",
            "始终使用中文回答。",
        ],
        "has_knowledge": True,
    },
    "a3": {
        "name": "代码助手Agent",
        "avatar": "💻",
        "description": "代码生成、审查、重构与调试，可执行 Python 代码",
        "capabilities": ["Python 执行", "代码生成", "文件读写", "调试"],
        "builtin_tools": ["python", "file"],
        "instructions": [
            "你是专业代码助手，擅长代码生成、审查、重构与调试。",
            "你可以直接执行 Python 代码并返回结果，也可以读写文件。",
            "提供高质量、可运行的代码示例，并用中文解释关键逻辑。",
            "始终使用中文回答。",
        ],
    },
    "a4": {
        "name": "合同审查Agent",
        "avatar": "📄",
        "description": "法律条款风险识别与比对",
        "capabilities": ["NLP", "合规检查", "文件读取"],
        "builtin_tools": ["file"],
        "instructions": [
            "你是法律合同审查专家，擅长识别合同风险和条款分析。",
            "你可以读取文件进行分析。",
            "指出潜在风险点，并给出修改建议。",
            "始终使用中文回答。",
        ],
    },
    "a5": {
        "name": "舆情监控Agent",
        "avatar": "📡",
        "description": "全网品牌舆情实时抓取与分析",
        "capabilities": ["情感分析", "趋势研判", "告警"],
        "instructions": [
            "你是品牌舆情分析专家，擅长情感分析和舆情趋势研判。",
            "始终使用中文回答。",
        ],
    },
    "a6": {
        "name": "私有数据治理Agent",
        "avatar": "🛡️",
        "description": "企业数据质量管理与 ETL 流程优化",
        "capabilities": ["数据质量", "规则引擎", "ETL"],
        "instructions": [
            "你是企业数据治理专家，擅长数据质量管理和 ETL 流程优化。",
            "始终使用中文回答。",
        ],
    },
    "a7": {
        "name": "PPT制作Agent",
        "avatar": "📑",
        "description": "根据描述自动生成专业 PPT 演示文稿",
        "capabilities": ["PPT生成", "模板设计", "内容编排"],
        "builtin_tools": ["python", "file"],
        "instructions": [
            "你是专业的 PPT 制作专家，能根据用户需求生成高质量的演示文稿。",
            "你使用 python-pptx 库来生成 .pptx 文件。",
            "生成 PPT 时请遵循以下原则：",
            "1. 首页包含标题和副标题",
            "2. 每页幻灯片有清晰的标题和要点",
            "3. 合理使用布局（标题页、内容页、双栏页等）",
            "4. 内容简洁，每页不超过 5-6 个要点",
            "5. 文件保存到工作目录",
            "",
            "生成 PPT 的代码模板：",
            "```python",
            "from pptx import Presentation",
            "from pptx.util import Inches, Pt",
            "from pptx.enum.text import PP_ALIGN",
            "prs = Presentation()",
            "# 标题页",
            "slide = prs.slides.add_slide(prs.slide_layouts[0])",
            "slide.shapes.title.text = '标题'",
            "slide.placeholders[1].text = '副标题'",
            "# 内容页",
            "slide = prs.slides.add_slide(prs.slide_layouts[1])",
            "slide.shapes.title.text = '章节标题'",
            "body = slide.placeholders[1]",
            "body.text = '要点内容'",
            "prs.save('output.pptx')",
            "```",
            "",
            "始终使用中文回答。生成完文件后，告知用户文件名。",
        ],
    },
    "global": {
        "name": "全局助手",
        "avatar": "🌐",
        "description": "效能管理智能体，管理项目、Agent 和 Skill 资源",
        "capabilities": ["资源查询", "效能分析", "计算器"],
        "builtin_tools": ["calculator"],
        "instructions": [
            "你是「全局助手」，帮助用户查询和管理系统中的 Agent、技能、项目和知识库资源。",
            "你有以下工具可用：",
            "- _global_list_agents：列出所有 Agent",
            "- _global_list_skills：列出所有技能",
            "- _global_list_projects：列出所有项目",
            "- _global_list_knowledge_docs：列出知识库文档",
            "- _global_system_stats：获取系统统计概览",
            "当用户询问 Agent、技能、项目、知识库相关信息时，主动调用对应工具获取实时数据，不要凭空编造。",
            "回答要简洁实用，使用列表或表格格式呈现结构化数据。",
            "始终使用中文回答。",
        ],
    },
    "skill_engineer": {
        "name": "技能工程师",
        "avatar": "🔧",
        "description": "根据自然语言描述自动生成 Python 技能",
        "capabilities": ["代码生成"],
        "instructions": [
            "你是技能工程师，专门根据用户的自然语言描述生成 Python 技能代码。",
            "生成的代码必须严格遵循以下格式：",
            "",
            "```python",
            'SKILL_META = {',
            '    "name": "技能中文名",',
            '    "icon": "合适的emoji",',
            '    "category": "analysis|data|code|search|api 之一",',
            '    "description": "一句话描述技能功能",',
            '}',
            "",
            "def run(参数1: 类型, 参数2: 类型) -> str:",
            '    """函数文档"""',
            "    # 实现逻辑",
            '    return "结果字符串"',
            "```",
            "",
            "你的回答必须是一个 JSON 对象，格式为：",
            '{"filename": "snake_case_name.py", "code": "完整的Python代码"}',
            "只返回 JSON，不要包含任何其他文字或 markdown 标记。",
            "文件名使用英文 snake_case，不含中文。",
            "run() 函数的参数使用基础类型（str, int, float, bool），返回 str。",
        ],
    },
}


# ═══════════════════════════════════════════════════════════════════════════════
# 全局助手 Tool Functions（查询系统资源）
# ═══════════════════════════════════════════════════════════════════════════════

def _global_list_agents() -> str:
    """列出系统中所有可用的 Agent 及其能力描述。"""
    lines = []
    for aid, cfg in AGENT_CONFIGS.items():
        if aid == "skill_engineer":
            continue
        tools = cfg.get("builtin_tools", [])
        custom = [sid for sid in _agent_tools.get(aid, []) if sid in _skill_registry]
        tool_str = ", ".join(tools + custom) if (tools or custom) else "无"
        lines.append(f"- **{cfg['name']}** (ID: {aid}): {cfg.get('description', '')} | 工具: {tool_str}")
    return f"当前系统共有 {len(lines)} 个 Agent：\n" + "\n".join(lines)


def _global_list_skills() -> str:
    """列出系统中所有已注册的技能。"""
    scan_skills()
    if not _skill_registry:
        return "当前没有已注册的技能。"
    lines = []
    for s in _skill_registry.values():
        meta = s["meta"]
        mounted = [AGENT_CONFIGS[aid]["name"] for aid, sids in _agent_tools.items() if s["id"] in sids and aid in AGENT_CONFIGS]
        mount_str = ", ".join(mounted) if mounted else "未挂载"
        lines.append(f"- **{meta['name']}** (ID: {s['id']}): {meta.get('description', '')} | 挂载: {mount_str}")
    return f"当前系统共有 {len(lines)} 个技能：\n" + "\n".join(lines)


def _global_list_projects() -> str:
    """列出系统中所有项目。"""
    conn = _get_projects_conn()
    rows = conn.execute("SELECT id, name, description, status, created_at FROM projects ORDER BY created_at DESC").fetchall()
    if not rows:
        return "当前没有项目。"
    lines = []
    for r in rows:
        lines.append(f"- **{r[1]}** (ID: {r[0]}): {r[2]} | 状态: {r[3]}")
    return f"当前系统共有 {len(lines)} 个项目：\n" + "\n".join(lines)


def _global_list_knowledge_docs() -> str:
    """列出知识库中已上传的文档。"""
    docs = list_documents()
    if not docs:
        return "知识库当前没有文档。"
    lines = [f"- {d['doc_name']} ({d['chunks']} 个段落)" for d in docs]
    return f"知识库共有 {len(lines)} 个文档：\n" + "\n".join(lines)


def _global_system_stats() -> str:
    """获取系统整体统计概览。"""
    scan_skills()
    agent_count = sum(1 for k in AGENT_CONFIGS if k != "skill_engineer")
    skill_count = len(_skill_registry)
    docs = list_documents()
    workspace_files = list(WORKSPACE_DIR.glob("*")) if WORKSPACE_DIR.exists() else []
    return (
        f"系统概览：\n"
        f"- Agent 数量: {agent_count}\n"
        f"- 技能数量: {skill_count}\n"
        f"- 知识库文档: {len(docs)}\n"
        f"- 工作区文件: {len(workspace_files)}"
    )


_GLOBAL_TOOLS = [
    _global_list_agents,
    _global_list_skills,
    _global_list_projects,
    _global_list_knowledge_docs,
    _global_system_stats,
]


# ═══════════════════════════════════════════════════════════════════════════════
# Agent 实例管理
# ═══════════════════════════════════════════════════════════════════════════════

_agents: dict[str, Agent] = {}


def _build_skill_tools(agent_id: str) -> list:
    """为指定 Agent 构建已挂载技能的工具函数列表。"""
    tool_ids = _agent_tools.get(agent_id, [])
    tools = []
    for sid in tool_ids:
        skill = _skill_registry.get(sid)
        if skill:
            tools.append(skill["run_fn"])
    return tools


def get_agent(agent_id: str) -> Agent | None:
    if agent_id in _agents:
        return _agents[agent_id]

    config = AGENT_CONFIGS.get(agent_id)
    if not config:
        return None

    try:
        builtin = _make_builtin_tools(config.get("builtin_tools", []))
        custom = _build_skill_tools(agent_id)
        all_tools = builtin + custom

        if agent_id == "global":
            all_tools = all_tools + _GLOBAL_TOOLS

        kwargs = {}
        if config.get("has_knowledge"):
            kwargs["knowledge"] = _knowledge
            kwargs["add_knowledge_to_context"] = True
            kwargs["search_knowledge"] = True
            kwargs["enable_agentic_knowledge_filters"] = True

        _guardrail = PromptInjectionGuardrail()

        agent = Agent(
            name=config["name"],
            id=agent_id,
            model=create_model(),
            db=SqliteDb(db_file=str(BASE_DIR / "data" / "sessions.db")),
            instructions=config["instructions"],
            tools=all_tools if all_tools else None,
            pre_hooks=[_guardrail],
            add_history_to_context=True,
            num_history_runs=10,
            markdown=True,
            **kwargs,
        )
        _agents[agent_id] = agent
        return agent
    except Exception as e:
        print(f"[ERROR] 创建 Agent {agent_id} 失败: {e}")
        return None


def invalidate_agent(agent_id: str):
    """清除 Agent 缓存，下次使用时重建（加载新工具）。"""
    _agents.pop(agent_id, None)
    _teams.clear()


# ═══════════════════════════════════════════════════════════════════════════════
# Team 多 Agent 协作 (coordinate 模式)
# ═══════════════════════════════════════════════════════════════════════════════

from agno.team.team import Team
from agno.team.mode import TeamMode

_teams: dict[str, Team] = {}

TEAM_MEMBER_IDS = ["a1", "a2", "a3"]

_INTENT_CATEGORIES: dict[str, dict] = {
    "knowledge":     {"target": "a2", "name": "知识检索Agent", "desc": "从知识库检索已有文档内容"},
    "data_analysis": {"target": "a1", "name": "数据分析Agent", "desc": "对 CSV/结构化数据做统计分析"},
    "code":          {"target": "a3", "name": "代码助手Agent", "desc": "代码生成、执行与调试"},
    "contract":      {"target": "a4", "name": "合同审查Agent", "desc": "法律合同条款风险识别"},
    "sentiment":     {"target": "a5", "name": "舆情监控Agent", "desc": "品牌舆情分析"},
    "data_govern":   {"target": "a6", "name": "私有数据治理Agent", "desc": "数据质量与 ETL 管理"},
    "ppt":           {"target": "a7", "name": "PPT制作Agent",  "desc": "生成演示文稿"},
    "multi_step":    {"target": None, "name": None, "desc": "需要多个 Agent 协作"},
}


async def classify_intent(message: str, docs: list[dict]) -> str:
    """用一次轻量 LLM 调用对用户问题做意图分类，返回类别标签。"""
    doc_list = ", ".join(d["doc_name"] for d in docs) if docs else "（无文档）"

    prompt = (
        "你是意图分类器。根据用户问题和已有知识库文档列表，判断该问题最适合由哪个类别处理。\n"
        "只返回一个分类标签（英文小写），不要解释。\n\n"
        f"已有知识库文档：{doc_list}\n"
        f"用户问题：{message}\n\n"
        "分类标签：\n"
        "- knowledge: 与已有文档/报告/简报/报表内容相关的问答检索（只要问题可能涉及已有文档就选此项）\n"
        "- data_analysis: 明确要求对 CSV/数据库做统计分析、建模、可视化\n"
        "- code: 代码生成、调试、审查\n"
        "- contract: 合同/法律条款审查\n"
        "- sentiment: 品牌舆情分析\n"
        "- data_govern: 数据质量/ETL 治理\n"
        "- ppt: 生成演示文稿/PPT\n"
        "- multi_step: 需要多个Agent协作的复杂任务（如先检索再分析再生成）\n"
    )

    try:
        classifier = Agent(
            model=create_model(),
            instructions=["只返回一个英文分类标签，不要输出任何其他内容。"],
            markdown=False,
        )
        response = await classifier.arun(prompt)
        label = (response.content or "").strip().lower().replace('"', "").replace("'", "")
        for key in _INTENT_CATEGORIES:
            if key in label:
                return key
        return "knowledge"
    except Exception as e:
        print(f"[WARN] 意图分类失败，fallback 到 knowledge: {e}")
        return "knowledge"


def get_team(project_id: str) -> Team:
    if project_id in _teams:
        return _teams[project_id]

    members = [get_agent(aid) for aid in TEAM_MEMBER_IDS]
    members = [m for m in members if m is not None]

    team = Team(
        name="项目团队",
        mode=TeamMode.route,
        model=create_model(),
        members=members,
        determine_input_for_members=False,
        db=SqliteDb(db_file=str(BASE_DIR / "data" / "sessions.db")),
        add_history_to_context=True,
        num_history_runs=5,
        add_team_history_to_members=True,
        num_team_history_runs=3,
        share_member_interactions=True,
        instructions=[
            "你是智能路由器，根据用户问题选择最合适的专家 Agent 来处理。",
            "你能看到之前的对话历史，要理解用户的指代（如'上面的'、'刚才的'等），结合上下文做出正确路由。",
            "",
            "## 核心路由原则",
            "当用户问的是文档、简报、报告中已有的内容（如排名、统计数据、分析结论等），必须选择知识检索Agent（a2），而不是数据分析Agent。",
            "只有当用户明确要求对 CSV 文件或数据库进行 SQL 查询、建模、可视化时，才选择数据分析Agent（a1）。",
            "当用户要求基于之前对话结果做进一步操作（如生成PPT、写报告），根据目标操作选择对应Agent，并确保上下文传递。",
            "",
            "## 成员能力",
            "- a2 知识检索Agent：从已上传的文档/报告/简报中检索信息，包括其中的排名、数据、分析结论等",
            "- a1 数据分析Agent：对 CSV/数据库执行 SQL 查询、Pandas 数据处理、统计建模",
            "- a3 代码助手Agent：代码生成、执行与调试",
            "- a4 合同审查Agent：法律合同条款风险识别",
            "- a5 舆情监控Agent：品牌舆情分析",
            "- a6 私有数据治理Agent：数据质量管理与 ETL",
            "- a7 PPT制作Agent：自动生成演示文稿",
            "",
            "## 路由规则",
            "- 涉及文档/报告/简报中的内容（排名、分析、数据、结论等）-> 知识检索Agent (a2)",
            "- 明确要求对 CSV/数据库做 SQL 查询、建模、可视化 -> 数据分析Agent (a1)",
            "- 代码生成/调试/审查 -> 代码助手Agent (a3)",
            "- 合同/法律相关 -> 合同审查Agent (a4)",
            "- 舆情/品牌相关 -> 舆情监控Agent (a5)",
            "- 数据治理/ETL -> 私有数据治理Agent (a6)",
            "- 制作PPT/演示文稿 -> PPT制作Agent (a7)",
            "- 通用技术方案/架构设计等 -> 代码助手Agent (a3)",
            "- 不确定时优先选择知识检索Agent (a2)",
        ],
        show_members_responses=True,
        markdown=True,
    )
    _teams[project_id] = team
    return team


# ═══════════════════════════════════════════════════════════════════════════════
# 请求模型
# ═══════════════════════════════════════════════════════════════════════════════

class ChatRequest(BaseModel):
    message: str
    session_id: str = "default"

class SkillCreateRequest(BaseModel):
    description: str

class SkillRunRequest(BaseModel):
    params: dict

class AgentToolsRequest(BaseModel):
    skill_ids: list[str]


# ═══════════════════════════════════════════════════════════════════════════════
# 路由：LLM 配置管理
# ═══════════════════════════════════════════════════════════════════════════════

def _mask_api_key(key: str) -> str:
    """脱敏 API Key，保留前 6 位和后 4 位。"""
    if not key or len(key) <= 12:
        return key
    return key[:6] + "..." + key[-4:]


class LLMSettingsRequest(BaseModel):
    provider: str
    model_id: str = ""
    api_key: str = ""
    base_url: str = ""


@app.get("/api/settings/llm")
async def api_get_llm_settings():
    cfg = _get_llm_config()
    return {
        "provider": cfg["provider"],
        "model_id": cfg["model_id"],
        "api_key": _mask_api_key(cfg["api_key"]),
        "base_url": cfg["base_url"],
        "providers": list(_PROVIDER_DEFAULT_MODEL.keys()),
        "default_models": _PROVIDER_DEFAULT_MODEL,
    }


@app.put("/api/settings/llm")
async def api_save_llm_settings(req: LLMSettingsRequest):
    now = datetime.now(timezone.utc).isoformat()
    conn = _get_projects_conn()

    api_key = req.api_key
    if "..." in api_key:
        existing = conn.execute(
            "SELECT value FROM settings WHERE key = 'llm_api_key'"
        ).fetchone()
        if existing:
            api_key = existing["value"]
        else:
            cfg = _get_llm_config()
            api_key = cfg["api_key"]

    pairs = {
        "llm_provider": req.provider,
        "llm_model_id": req.model_id,
        "llm_api_key": api_key,
        "llm_base_url": req.base_url,
    }
    for k, v in pairs.items():
        conn.execute(
            "INSERT INTO settings (key, value, updated_at) VALUES (?, ?, ?) "
            "ON CONFLICT(key) DO UPDATE SET value = excluded.value, updated_at = excluded.updated_at",
            (k, v, now),
        )
    conn.commit()
    conn.close()

    _agents.clear()
    _teams.clear()

    return {"success": True}


@app.post("/api/settings/llm/test")
async def api_test_llm_connection(req: LLMSettingsRequest):
    """用提交的配置做一次轻量 LLM 调用，测试连通性。"""
    api_key = req.api_key
    if "..." in api_key:
        conn = _get_projects_conn()
        existing = conn.execute(
            "SELECT value FROM settings WHERE key = 'llm_api_key'"
        ).fetchone()
        conn.close()
        if existing:
            api_key = existing["value"]
        else:
            cfg = _get_llm_config()
            api_key = cfg["api_key"]

    try:
        from agno.models.openai import OpenAIChat

        provider = req.provider.lower()
        model_id = req.model_id or _PROVIDER_DEFAULT_MODEL.get(provider, "gpt-4o-mini")
        base_url = req.base_url or _PROVIDER_BASE_URL_MAP.get(provider)

        kwargs = dict(id=model_id, timeout=15)
        if api_key:
            kwargs["api_key"] = api_key
        if base_url:
            kwargs["base_url"] = base_url
        kwargs.update(_PROVIDER_EXTRA_KWARGS.get(provider, {}))
        model = OpenAIChat(**kwargs)

        test_agent = Agent(model=model, instructions=["回复OK"], markdown=False)
        resp = await test_agent.arun("ping", stream=False)
        content = (resp.content or "").strip()
        error_keywords = ["invalid", "unauthorized", "error", "failed", "denied", "forbidden", "401", "403"]
        content_lower = content.lower()
        if any(kw in content_lower for kw in error_keywords) or not content:
            return {"ok": False, "message": f"认证失败，模型返回: {content[:200]}"}
        return {"ok": True, "message": f"连接成功，模型响应: {content[:100]}"}
    except Exception as e:
        return {"ok": False, "message": f"连接失败: {str(e)[:200]}"}


# ═══════════════════════════════════════════════════════════════════════════════
# 路由：Embedding 配置
# ═══════════════════════════════════════════════════════════════════════════════

class EmbeddingSettingsRequest(BaseModel):
    mode: str = "local"
    model_id: str = ""
    api_key: str = ""
    base_url: str = ""
    dimensions: int = 1024


def _mask_key(key: str) -> str:
    if not key or len(key) <= 8:
        return key
    return key[:4] + "..." + key[-4:]


@app.get("/api/settings/embedding")
async def api_get_embedding_settings():
    cfg = _get_embedding_config()
    cfg["api_key"] = _mask_key(cfg["api_key"])
    return cfg


@app.put("/api/settings/embedding")
async def api_save_embedding_settings(req: EmbeddingSettingsRequest):
    now = datetime.now(timezone.utc).isoformat()
    conn = _get_projects_conn()

    api_key = req.api_key
    if "..." in api_key:
        existing = conn.execute(
            "SELECT value FROM settings WHERE key = 'embedding_api_key'"
        ).fetchone()
        if existing:
            api_key = existing["value"]
        else:
            api_key = _get_embedding_config()["api_key"]

    if req.mode == "local":
        for k in ["embedding_model_id", "embedding_api_key", "embedding_base_url", "embedding_dimensions"]:
            conn.execute("DELETE FROM settings WHERE key = ?", (k,))
    else:
        pairs = {
            "embedding_model_id": req.model_id,
            "embedding_api_key": api_key,
            "embedding_base_url": req.base_url,
            "embedding_dimensions": str(req.dimensions),
        }
        for k, v in pairs.items():
            conn.execute(
                "INSERT INTO settings (key, value, updated_at) VALUES (?, ?, ?) "
                "ON CONFLICT(key) DO UPDATE SET value = excluded.value, updated_at = excluded.updated_at",
                (k, v, now),
            )
    conn.commit()
    conn.close()

    _rebuild_knowledge()
    _agents.clear()
    _teams.clear()
    return {"success": True, "warning": "切换 Embedding 模型后，建议重新上传文档以重建索引。"}


@app.post("/api/settings/embedding/test")
async def api_test_embedding_connection(req: EmbeddingSettingsRequest):
    """调用 /v1/embeddings 测试 Embedding 模型连通性。"""
    api_key = req.api_key
    if "..." in api_key:
        api_key = _get_embedding_config()["api_key"]

    if req.mode == "local":
        return {"ok": True, "message": "本地 FastEmbed 模式，无需测试连通性。"}

    try:
        from openai import OpenAI
        client = OpenAI(
            api_key=api_key or "none",
            base_url=req.base_url,
            timeout=15,
        )
        resp = client.embeddings.create(
            model=req.model_id,
            input=["测试连通性"],
        )
        dim = len(resp.data[0].embedding)
        return {"ok": True, "message": f"连接成功，返回向量维度: {dim}"}
    except Exception as e:
        return {"ok": False, "message": f"连接失败: {str(e)[:200]}"}


# ═══════════════════════════════════════════════════════════════════════════════
# 路由：Reranker 配置
# ═══════════════════════════════════════════════════════════════════════════════

class RerankerSettingsRequest(BaseModel):
    enabled: bool = False
    model_id: str = ""
    api_key: str = ""
    base_url: str = ""
    top_n: int = 5


@app.get("/api/settings/reranker")
async def api_get_reranker_settings():
    cfg = _get_reranker_config()
    cfg["api_key"] = _mask_key(cfg["api_key"])
    return cfg


@app.put("/api/settings/reranker")
async def api_save_reranker_settings(req: RerankerSettingsRequest):
    now = datetime.now(timezone.utc).isoformat()
    conn = _get_projects_conn()

    api_key = req.api_key
    if "..." in api_key:
        existing = conn.execute(
            "SELECT value FROM settings WHERE key = 'reranker_api_key'"
        ).fetchone()
        if existing:
            api_key = existing["value"]
        else:
            api_key = _get_reranker_config()["api_key"]

    pairs = {
        "reranker_enabled": str(req.enabled).lower(),
        "reranker_model_id": req.model_id,
        "reranker_api_key": api_key,
        "reranker_base_url": req.base_url,
        "reranker_top_n": str(req.top_n),
    }
    for k, v in pairs.items():
        conn.execute(
            "INSERT INTO settings (key, value, updated_at) VALUES (?, ?, ?) "
            "ON CONFLICT(key) DO UPDATE SET value = excluded.value, updated_at = excluded.updated_at",
            (k, v, now),
        )
    conn.commit()
    conn.close()

    _rebuild_knowledge()
    _agents.clear()
    _teams.clear()
    return {"success": True}


@app.post("/api/settings/reranker/test")
async def api_test_reranker_connection(req: RerankerSettingsRequest):
    """调用 rerank API 测试 Reranker 连通性。"""
    api_key = req.api_key
    if "..." in api_key:
        api_key = _get_reranker_config()["api_key"]

    if not req.enabled:
        return {"ok": True, "message": "Reranker 未启用，无需测试。"}

    try:
        reranker = OpenAICompatibleReranker(
            model=req.model_id,
            api_key=api_key,
            base_url=req.base_url,
            top_n=req.top_n,
        )
        test_docs = ["人工智能是计算机科学的分支", "今天天气很好适合出去玩", "深度学习需要大量训练数据"]
        results = reranker._call_rerank("什么是人工智能", test_docs)
        if results:
            top = results[0]
            return {"ok": True, "message": f"连接成功，返回 {len(results)} 条排序结果，最高分: {top.get('relevance_score', 'N/A')}"}
        return {"ok": False, "message": "连接成功但返回结果为空"}
    except Exception as e:
        return {"ok": False, "message": f"连接失败: {str(e)[:200]}"}


# ═══════════════════════════════════════════════════════════════════════════════
# 路由：健康检查 & Agent 聊天
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/health")
async def health():
    cfg = _get_llm_config()
    return {
        "status": "ok",
        "model_provider": cfg["provider"],
        "model_id": cfg["model_id"],
        "skills_count": len(_skill_registry),
        "docs_count": len(list_documents()),
    }


# ═══════════════════════════════════════════════════════════════════════════════
# 路由：项目管理
# ═══════════════════════════════════════════════════════════════════════════════

class ProjectCreateRequest(BaseModel):
    name: str
    description: str = ""


@app.get("/api/projects")
async def api_list_projects():
    conn = _get_projects_conn()
    rows = conn.execute("SELECT * FROM projects ORDER BY updated_at DESC").fetchall()
    conn.close()
    return [dict(r) for r in rows]


@app.post("/api/projects")
async def api_create_project(request: ProjectCreateRequest):
    import uuid
    project_id = f"p{uuid.uuid4().hex[:8]}"
    now = datetime.now(timezone.utc).isoformat()
    conn = _get_projects_conn()
    conn.execute(
        "INSERT INTO projects (id, name, description, status, created_at, updated_at) VALUES (?,?,?,?,?,?)",
        (project_id, request.name, request.description, "active", now, now),
    )
    conn.commit()
    row = conn.execute("SELECT * FROM projects WHERE id = ?", (project_id,)).fetchone()
    conn.close()
    return dict(row)


@app.delete("/api/projects/{project_id}")
async def api_delete_project(project_id: str):
    conn = _get_projects_conn()
    conn.execute("DELETE FROM tasks WHERE project_id = ?", (project_id,))
    conn.execute("DELETE FROM projects WHERE id = ?", (project_id,))
    conn.commit()
    conn.close()
    return {"success": True}


# ── Task CRUD ─────────────────────────────────────────────────────────────────

@app.get("/api/projects/{project_id}/tasks")
async def api_list_tasks(project_id: str):
    conn = _get_projects_conn()
    rows = conn.execute(
        "SELECT * FROM tasks WHERE project_id = ? ORDER BY sort_order, created_at",
        (project_id,),
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


class TaskCreateRequest(BaseModel):
    name: str


@app.post("/api/projects/{project_id}/tasks")
async def api_create_task(project_id: str, request: TaskCreateRequest):
    import uuid
    task_id = f"t{uuid.uuid4().hex[:8]}"
    now_str = datetime.now(timezone.utc).isoformat()
    conn = _get_projects_conn()
    max_order = conn.execute(
        "SELECT COALESCE(MAX(sort_order), 0) FROM tasks WHERE project_id = ?",
        (project_id,),
    ).fetchone()[0]
    conn.execute(
        "INSERT INTO tasks (id, project_id, name, sort_order, created_at) VALUES (?,?,?,?,?)",
        (task_id, project_id, request.name, max_order + 1, now_str),
    )
    conn.commit()
    row = conn.execute("SELECT * FROM tasks WHERE id = ?", (task_id,)).fetchone()
    conn.close()
    return dict(row)


@app.delete("/api/tasks/{task_id}")
async def api_delete_task(task_id: str):
    conn = _get_projects_conn()
    conn.execute("DELETE FROM tasks WHERE id = ?", (task_id,))
    conn.commit()
    conn.close()
    return {"success": True}


_HIDDEN_AGENT_IDS = {"skill_engineer", "global", "a4", "a5", "a6", "a7"}


@app.get("/api/agents")
async def api_list_agents():
    return [
        {
            "id": k,
            "name": v["name"],
            "avatar": v.get("avatar", "🤖"),
            "description": v.get("description", ""),
            "capabilities": v.get("capabilities", []),
            "builtin_tools": v.get("builtin_tools", []),
            "custom_tools": _agent_tools.get(k, []),
            "has_knowledge": v.get("has_knowledge", False),
            "instructions": v.get("instructions", []),
        }
        for k, v in AGENT_CONFIGS.items()
        if k not in _HIDDEN_AGENT_IDS
    ]


@app.post("/api/agents/{agent_id}/chat")
async def agent_chat(agent_id: str, request: ChatRequest):
    agent = get_agent(agent_id)

    async def generate():
        if not agent:
            yield f"data: {json.dumps({'content': f'Agent [{agent_id}] 未找到。', 'done': True})}\n\n"
            return
        try:
            async with asyncio.timeout(180):
                async for chunk in agent.arun(
                    request.message,
                    stream=True,
                    session_id=request.session_id,
                ):
                    if chunk.content:
                        cleaned = _clean_delta(chunk.content)
                        if cleaned:
                            yield f"data: {json.dumps({'content': cleaned, 'done': False}, ensure_ascii=False)}\n\n"
            yield f"data: {json.dumps({'content': '', 'done': True})}\n\n"
        except TimeoutError:
            print(f"[ERROR] Agent chat {agent_id}: timeout after 180s")
            yield f"data: {json.dumps({'content': 'Agent 响应超时（3分钟），请简化问题后重试。', 'done': True})}\n\n"
        except Exception as e:
            err = str(e)
            if "api_key" in err.lower() or "authentication" in err.lower():
                msg = "⚠️ 未检测到有效的 API Key，请在 backend/.env 中配置。"
            else:
                msg = f"请求出错：{err}"
            yield f"data: {json.dumps({'content': msg, 'done': True})}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


# ═══════════════════════════════════════════════════════════════════════════════
# 路由：Team 协作
# ═══════════════════════════════════════════════════════════════════════════════

_THINK_RE = re.compile(r"<think>.*?</think>", re.DOTALL)
_MEMBER_XML_RE = re.compile(r"<member\b[^>]*>.*?</member>", re.DOTALL)
_MEMBER_LIST_BLOCK_RE = re.compile(
    r"Please choose the correct member from the list of members:\s*(?:<member\b.*?</member>\s*)+",
    re.DOTALL,
)

def _clean_content(text: str) -> str:
    """清理完整消息中的 think 标签、成员列表 XML 等。"""
    text = _THINK_RE.sub("", text)
    text = _MEMBER_LIST_BLOCK_RE.sub("", text)
    text = _MEMBER_XML_RE.sub("", text)
    text = text.strip()
    return text

def _clean_delta(text: str) -> str:
    """清理流式 chunk 中的 think 标签，但保留空白和换行。"""
    return _THINK_RE.sub("", text)

def _resolve_agent_display(member_id: str | None) -> tuple[str, str]:
    """根据 member_id 解析出 (avatar, display_name)。"""
    if not member_id:
        return "🤖", "成员"
    for k, v in AGENT_CONFIGS.items():
        if k == member_id or v["name"] == member_id:
            return v.get("avatar", "🤖"), v["name"]
    return "🤖", member_id


@app.post("/api/teams/{project_id}/chat")
async def team_chat(project_id: str, request: ChatRequest):
    from agno.run.team import TeamRunEvent, TeamRunOutput
    from agno.run.agent import RunEvent, RunOutput

    team = get_team(project_id)
    streamed_members: set[str] = set()

    def _sse(payload: dict) -> str:
        return f"data: {json.dumps(payload, ensure_ascii=False)}\n\n"

    async def generate():
        routed_agent_name: str | None = None

        try:
            async with asyncio.timeout(300):
                async for event in team.arun(
                    request.message,
                    stream=True,
                    stream_events=True,
                    session_id=f"team_{project_id}_{request.session_id}",
                    yield_run_output=True,
                ):
                    if isinstance(event, TeamRunOutput):
                        for mr in event.member_responses:
                            mid = None
                            if isinstance(mr, RunOutput):
                                mid = mr.agent_id
                            if mid and mid in streamed_members:
                                continue
                            avatar, display_name = _resolve_agent_display(mid)
                            content = ""
                            if hasattr(mr, "content") and mr.content:
                                content = mr.content if isinstance(mr.content, str) else str(mr.content)
                            content = _clean_content(content)
                            if content:
                                yield _sse({
                                    "type": "member_response",
                                    "agent_name": f"{avatar} {display_name}",
                                    "content": content,
                                    "done": False,
                                })
                        continue

                    ev_type = getattr(event, "event", "")

                    if ev_type == TeamRunEvent.tool_call_started.value:
                        tool = getattr(event, "tool", None)
                        if tool and tool.tool_name and ("delegate_task" in tool.tool_name or "route_to" in tool.tool_name):
                            args = tool.tool_args or {}
                            mid = args.get("member_id", "")
                            task_raw = args.get("task", args.get("input", ""))
                            avatar, display_name = _resolve_agent_display(mid)
                            routed_agent_name = f"{avatar} {display_name}"
                            streamed_members.add(mid)
                            task_clean = (task_raw or "").replace("\\n", " ").strip()
                            if len(task_clean) > 80:
                                task_clean = task_clean[:80] + "..."
                            yield _sse({
                                "type": "member_delegated",
                                "agent_name": routed_agent_name,
                                "task": task_clean or request.message,
                                "done": False,
                            })

                    elif ev_type == TeamRunEvent.tool_call_completed.value:
                        pass

                    elif ev_type == TeamRunEvent.run_content.value:
                        content = getattr(event, "content", None)
                        if content:
                            text = content if isinstance(content, str) else str(content)
                            text = _clean_delta(text)
                            if text:
                                if routed_agent_name:
                                    yield _sse({
                                        "type": "member_streaming",
                                        "agent_name": routed_agent_name,
                                        "content": text,
                                        "done": False,
                                    })
                                else:
                                    yield _sse({
                                        "type": "leader_content",
                                        "content": text,
                                        "done": False,
                                    })

                    elif ev_type == RunEvent.run_content.value:
                        agent_id = getattr(event, "agent_id", "")
                        content = getattr(event, "content", None)
                        if content and agent_id:
                            streamed_members.add(agent_id)
                            avatar, display_name = _resolve_agent_display(agent_id)
                            text = content if isinstance(content, str) else str(content)
                            text = _clean_delta(text)
                            if text:
                                yield _sse({
                                    "type": "member_streaming",
                                    "agent_name": f"{avatar} {display_name}",
                                    "content": text,
                                    "done": False,
                                })

                    elif ev_type == RunEvent.run_started.value:
                        agent_id = getattr(event, "agent_id", "")
                        if agent_id:
                            avatar, display_name = _resolve_agent_display(agent_id)
                            yield _sse({
                                "type": "member_started",
                                "agent_name": f"{avatar} {display_name}",
                                "done": False,
                            })

            yield _sse({"type": "done", "content": "", "done": True})
        except TimeoutError:
            print("[ERROR] Team chat: timeout after 300s")
            yield _sse({"type": "leader_content", "content": "Team 协作超时（5分钟），请简化问题后重试。", "done": True})
        except Exception as e:
            err = str(e)
            print(f"[ERROR] Team chat: {err}")
            yield _sse({"type": "leader_content", "content": f"Team 协作出错：{err}", "done": True})

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


# ═══════════════════════════════════════════════════════════════════════════════
# 路由：技能管理
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/api/skills")
async def api_list_skills():
    scan_skills()
    return [
        {
            "id": s["id"],
            "name": s["meta"]["name"],
            "icon": s["meta"].get("icon", "🔧"),
            "category": s["meta"].get("category", "api"),
            "description": s["meta"].get("description", ""),
            "params": s["params"],
            "mounted_agents": [
                {"id": aid, "name": AGENT_CONFIGS[aid]["name"]}
                for aid, sids in _agent_tools.items()
                if s["id"] in sids and aid in AGENT_CONFIGS
            ],
        }
        for s in _skill_registry.values()
    ]


@app.post("/api/skills/create")
async def api_create_skill(request: SkillCreateRequest):
    """调用技能工程师 Agent 根据自然语言描述生成技能代码。"""
    engineer = get_agent("skill_engineer")
    if not engineer:
        return {"success": False, "error": "技能工程师 Agent 初始化失败"}

    try:
        result = await engineer.arun(request.description, stream=False)
        raw = result.content.strip()
        # 尝试从回复中提取 JSON
        if "```" in raw:
            match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', raw, re.DOTALL)
            if match:
                raw = match.group(1)
        data = json.loads(raw)
        filename = data["filename"]
        code = data["code"]

        if not filename.endswith(".py"):
            filename += ".py"
        filename = re.sub(r'[^a-zA-Z0-9_.]', '_', filename)

        filepath = SKILLS_DIR / filename
        filepath.write_text(code, encoding="utf-8")

        scan_skills()
        skill_id = filepath.stem
        if skill_id in _skill_registry:
            return {"success": True, "skill_id": skill_id, "skill": _skill_registry[skill_id]["meta"]}
        else:
            return {"success": False, "error": "技能代码生成成功但加载失败，请检查代码格式"}

    except json.JSONDecodeError:
        return {"success": False, "error": f"技能工程师返回了非 JSON 格式的内容: {raw[:200]}"}
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.post("/api/skills/{skill_id}/run")
async def api_run_skill(skill_id: str, request: SkillRunRequest):
    skill = _skill_registry.get(skill_id)
    if not skill:
        return {"success": False, "error": f"技能 [{skill_id}] 不存在"}
    try:
        sig = inspect.signature(skill["run_fn"])
        cast_params = {}
        for name, param in sig.parameters.items():
            if name in request.params:
                val = request.params[name]
                if param.annotation == float:
                    val = float(val)
                elif param.annotation == int:
                    val = int(val)
                cast_params[name] = val
        result = skill["run_fn"](**cast_params)
        return {"success": True, "result": str(result)}
    except Exception as e:
        return {"success": False, "error": str(e)}


class CreateAgentRequest(BaseModel):
    name: str
    avatar: str = "🤖"
    description: str = ""
    instructions: list[str] | None = None
    skill_ids: list[str] | None = None
    builtin_tools: list[str] | None = None
    join_team: bool = False


@app.post("/api/agents")
async def api_create_agent(req: CreateAgentRequest):
    agent_id = f"custom_{int(time.time() * 1000)}"
    config = {
        "name": req.name,
        "avatar": req.avatar,
        "description": req.description,
        "capabilities": [],
        "builtin_tools": req.builtin_tools or [],
        "instructions": req.instructions or [
            f"你是{req.name}，一个智能数字员工。",
            "根据用户的指令完成任务。",
            "始终使用中文回答。",
        ],
    }
    AGENT_CONFIGS[agent_id] = config
    if req.skill_ids:
        _agent_tools[agent_id] = req.skill_ids
    if req.join_team:
        TEAM_MEMBER_IDS.append(agent_id)
        _teams.clear()
    return {
        "success": True,
        "agent": {
            "id": agent_id,
            "name": config["name"],
            "avatar": config["avatar"],
            "description": config["description"],
            "capabilities": config["capabilities"],
            "builtin_tools": config["builtin_tools"],
            "custom_tools": _agent_tools.get(agent_id, []),
            "has_knowledge": config.get("has_knowledge", False),
            "instructions": config["instructions"],
        },
    }


@app.delete("/api/agents/{agent_id}")
async def api_delete_agent(agent_id: str):
    if not agent_id.startswith("custom_"):
        return {"success": False, "error": "内置 Agent 不允许删除"}
    if agent_id not in AGENT_CONFIGS:
        return {"success": False, "error": f"Agent [{agent_id}] 不存在"}
    AGENT_CONFIGS.pop(agent_id, None)
    _agent_tools.pop(agent_id, None)
    _agents.pop(agent_id, None)
    if agent_id in TEAM_MEMBER_IDS:
        TEAM_MEMBER_IDS.remove(agent_id)
    _teams.clear()
    return {"success": True, "agent_id": agent_id}


@app.put("/api/agents/{agent_id}/tools")
async def api_set_agent_tools(agent_id: str, request: AgentToolsRequest):
    if agent_id not in AGENT_CONFIGS:
        return {"success": False, "error": f"Agent [{agent_id}] 不存在"}
    _agent_tools[agent_id] = request.skill_ids
    invalidate_agent(agent_id)
    return {"success": True, "agent_id": agent_id, "tools": request.skill_ids}


@app.delete("/api/skills/{skill_id}")
async def api_delete_skill(skill_id: str):
    skill = _skill_registry.get(skill_id)
    if not skill:
        return {"success": False, "error": f"技能 [{skill_id}] 不存在"}
    filepath = SKILLS_DIR / f"{skill_id}.py"
    if filepath.exists():
        filepath.unlink()
    for aid, sids in _agent_tools.items():
        if skill_id in sids:
            sids.remove(skill_id)
            invalidate_agent(aid)
    _skill_registry.pop(skill_id, None)
    return {"success": True, "skill_id": skill_id}


class AgentConfigRequest(BaseModel):
    description: str | None = None
    instructions: list[str] | None = None


@app.put("/api/agents/{agent_id}/config")
async def api_update_agent_config(agent_id: str, request: AgentConfigRequest):
    if agent_id not in AGENT_CONFIGS:
        return {"success": False, "error": f"Agent [{agent_id}] 不存在"}
    config = AGENT_CONFIGS[agent_id]
    if request.description is not None:
        config["description"] = request.description
    if request.instructions is not None:
        config["instructions"] = request.instructions
    invalidate_agent(agent_id)
    return {"success": True, "agent_id": agent_id}


# ═══════════════════════════════════════════════════════════════════════════════
# 技能管理 Agent（对话式管理技能）
# ═══════════════════════════════════════════════════════════════════════════════

def _skill_tool_mount(skill_id: str, agent_id: str) -> str:
    """将技能挂载到指定 Agent。"""
    if agent_id not in AGENT_CONFIGS:
        return f"错误：Agent [{agent_id}] 不存在。可用 Agent：{', '.join(k + '(' + v['name'] + ')' for k, v in AGENT_CONFIGS.items() if k != 'skill_engineer')}"
    if skill_id not in _skill_registry:
        return f"错误：技能 [{skill_id}] 不存在"
    current = _agent_tools.get(agent_id, [])
    if skill_id in current:
        return f"技能 [{skill_id}] 已经挂载在 {AGENT_CONFIGS[agent_id]['name']} 上了"
    _agent_tools[agent_id] = current + [skill_id]
    invalidate_agent(agent_id)
    return f"成功：已将技能挂载到 {AGENT_CONFIGS[agent_id]['name']}（{agent_id}）"


def _skill_tool_unmount(skill_id: str, agent_id: str) -> str:
    """从指定 Agent 卸载技能。"""
    if agent_id not in AGENT_CONFIGS:
        return f"错误：Agent [{agent_id}] 不存在"
    current = _agent_tools.get(agent_id, [])
    if skill_id not in current:
        return f"技能 [{skill_id}] 未挂载在 {AGENT_CONFIGS[agent_id]['name']} 上"
    current.remove(skill_id)
    _agent_tools[agent_id] = current
    invalidate_agent(agent_id)
    return f"成功：已从 {AGENT_CONFIGS[agent_id]['name']} 卸载技能"


def _skill_tool_run(skill_id: str, params_json: str) -> str:
    """执行技能，params_json 是 JSON 格式的参数字典。"""
    skill = _skill_registry.get(skill_id)
    if not skill:
        return f"错误：技能 [{skill_id}] 不存在"
    try:
        params = json.loads(params_json) if params_json.strip() else {}
        sig = inspect.signature(skill["run_fn"])
        cast = {}
        for name, param in sig.parameters.items():
            if name in params:
                val = params[name]
                if param.annotation == float:
                    val = float(val)
                elif param.annotation == int:
                    val = int(val)
                cast[name] = val
        result = skill["run_fn"](**cast)
        return f"执行成功：{result}"
    except Exception as e:
        return f"执行失败：{e}"


def _skill_tool_delete(skill_id: str) -> str:
    """删除技能文件并清除所有 Agent 绑定。"""
    if skill_id not in _skill_registry:
        return f"错误：技能 [{skill_id}] 不存在"
    filepath = SKILLS_DIR / f"{skill_id}.py"
    if filepath.exists():
        filepath.unlink()
    for aid, sids in _agent_tools.items():
        if skill_id in sids:
            sids.remove(skill_id)
            invalidate_agent(aid)
    _skill_registry.pop(skill_id, None)
    return f"成功：技能 [{skill_id}] 已删除，相关 Agent 绑定已清除"


async def _skill_tool_modify(skill_id: str, instruction: str) -> str:
    """根据自然语言指令修改技能代码。读取当前代码，用 LLM 改写后替换。"""
    if skill_id not in _skill_registry:
        return f"错误：技能 [{skill_id}] 不存在"
    filepath = SKILLS_DIR / f"{skill_id}.py"
    if not filepath.exists():
        return f"错误：技能文件不存在"
    current_code = filepath.read_text(encoding="utf-8")

    engineer = get_agent("skill_engineer")
    if not engineer:
        return "错误：技能工程师 Agent 初始化失败"

    prompt = (
        f"请修改以下 Python 技能代码。修改要求：{instruction}\n\n"
        f"当前代码：\n```python\n{current_code}\n```\n\n"
        f"请返回一个 JSON 对象：{{\"filename\": \"{skill_id}.py\", \"code\": \"修改后的完整 Python 代码\"}}\n"
        f"只返回 JSON，不要包含任何其他文字或 markdown 标记。"
    )
    try:
        result = await engineer.arun(prompt, stream=False)
        raw = result.content.strip()
        if "```" in raw:
            match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', raw, re.DOTALL)
            if match:
                raw = match.group(1)
        data = json.loads(raw)
        new_code = data["code"]
        filepath.write_text(new_code, encoding="utf-8")
        scan_skills()
        if skill_id in _skill_registry:
            return f"成功：技能代码已更新。新描述：{_skill_registry[skill_id]['meta'].get('description', '')}"
        else:
            return "警告：代码已写入但加载失败，可能存在语法错误"
    except Exception as e:
        return f"修改失败：{e}"


async def _skill_tool_create(description: str) -> str:
    """根据自然语言描述创建新技能。"""
    engineer = get_agent("skill_engineer")
    if not engineer:
        return "错误：技能工程师 Agent 初始化失败"

    try:
        result = await engineer.arun(description, stream=False)
        raw = result.content.strip()
        if "```" in raw:
            match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', raw, re.DOTALL)
            if match:
                raw = match.group(1)
        data = json.loads(raw)
        filename = data["filename"]
        code = data["code"]
        if not filename.endswith(".py"):
            filename += ".py"
        filename = re.sub(r'[^a-zA-Z0-9_.]', '_', filename)
        filepath = SKILLS_DIR / filename
        filepath.write_text(code, encoding="utf-8")
        scan_skills()
        sid = filepath.stem
        if sid in _skill_registry:
            meta = _skill_registry[sid]["meta"]
            return f"成功：已创建技能「{meta['name']}」(ID: {sid})，描述：{meta.get('description', '')}"
        else:
            return "代码已生成但加载失败，请检查代码格式"
    except Exception as e:
        return f"创建失败：{e}"


def _build_skill_manager_agent(skill_id: str | None):
    """构建技能管理 Agent，注入当前技能上下文。"""
    context_lines = []

    if skill_id and skill_id in _skill_registry:
        s = _skill_registry[skill_id]
        meta = s["meta"]
        params_desc = ", ".join(f"{p['name']}:{p['type']}" for p in s["params"]) or "无参数"
        mounted = [
            f"{AGENT_CONFIGS[aid]['name']}({aid})"
            for aid, sids in _agent_tools.items()
            if skill_id in sids and aid in AGENT_CONFIGS
        ]
        context_lines.append(f"当前技能：{meta['name']}（ID: {skill_id}）")
        context_lines.append(f"描述：{meta.get('description', '')}")
        context_lines.append(f"分类：{meta.get('category', '')}")
        context_lines.append(f"参数：{params_desc}")
        context_lines.append(f"已挂载到：{', '.join(mounted) if mounted else '未挂载'}")
    else:
        context_lines.append("当前未选中特定技能，你可以帮用户创建新技能。")

    agent_list = ", ".join(
        f"{k}({v['name']})" for k, v in AGENT_CONFIGS.items() if k != "skill_engineer"
    )
    context_lines.append(f"\n可用 Agent 列表：{agent_list}")

    skill_list = ", ".join(
        f"{s['id']}({s['meta']['name']})" for s in _skill_registry.values()
    )
    if skill_list:
        context_lines.append(f"已注册技能列表：{skill_list}")

    context = "\n".join(context_lines)

    instructions = [
        "你是技能管理助手，帮助用户通过自然语言管理技能。",
        "你拥有以下能力：挂载/卸载技能到 Agent、执行技能、修改技能代码、创建新技能、删除技能。",
        "用户可能用模糊的说法，比如'挂到数据分析上'，你需要推断对应的 agent_id。",
        "执行技能时，从用户消息中提取参数值，构造 JSON 字符串传给 run 工具。",
        "修改技能时，将用户的修改意图转为清晰的指令传给 modify 工具。",
        "回复使用中文，简洁明了。",
        f"\n--- 上下文 ---\n{context}",
    ]

    tools = [
        _skill_tool_mount,
        _skill_tool_unmount,
        _skill_tool_run,
        _skill_tool_delete,
        _skill_tool_modify,
        _skill_tool_create,
    ]

    return Agent(
        model=create_model(),
        tools=tools,
        instructions=instructions,
        markdown=True,
    )


@app.post("/api/skills/{skill_id}/chat")
async def skill_chat(skill_id: str, request: ChatRequest):
    agent = _build_skill_manager_agent(skill_id)

    async def generate():
        try:
            async with asyncio.timeout(180):
                async for chunk in agent.arun(
                    request.message,
                    stream=True,
                    session_id=f"skill_mgr_{skill_id}_{request.session_id}",
                ):
                    if chunk.content:
                        cleaned = _clean_delta(chunk.content)
                        if cleaned:
                            yield f"data: {json.dumps({'content': cleaned, 'done': False}, ensure_ascii=False)}\n\n"
            yield f"data: {json.dumps({'content': '', 'done': True})}\n\n"
        except TimeoutError:
            yield f"data: {json.dumps({'content': '技能管理超时（3分钟），请简化操作后重试。', 'done': True})}\n\n"
        except Exception as e:
            err = str(e)
            print(f"[ERROR] Skill chat {skill_id}: {err}")
            yield f"data: {json.dumps({'content': f'出错了：{err}', 'done': True}, ensure_ascii=False)}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


# ═══════════════════════════════════════════════════════════════════════════════
# 路由：知识库
# ═══════════════════════════════════════════════════════════════════════════════

@app.post("/api/knowledge/upload")
async def api_upload_document(file: UploadFile = File(...)):
    content = await file.read()
    doc_name = file.filename or "unknown.txt"
    suffix = Path(doc_name).suffix.lower()

    doc_path = KNOWLEDGE_DOCS_DIR / doc_name

    if suffix == ".pdf":
        import fitz
        doc_path.write_bytes(content)
        pdf = fitz.open(stream=content, filetype="pdf")
        _end_puncts = set('。；！？.;!?\n')
        parts: list[str] = []
        for page in pdf:
            t = page.get_text().strip()
            if not t:
                continue
            if parts and parts[-1] and parts[-1][-1] not in _end_puncts:
                parts.append(' ' + t)
            else:
                parts.append('\n\n' + t if parts else t)
        pdf.close()
        text = ''.join(parts).strip()
    else:
        text = content.decode("utf-8", errors="ignore")
        doc_path.write_text(text, encoding="utf-8")

    if not text.strip():
        return {"success": False, "error": "文件内容为空，无法解析"}

    try:
        chunk_count = ingest_document(doc_name, text)
        _uploaded_docs[doc_name] = chunk_count
        invalidate_agent("a2")
        return {"success": True, "doc_name": doc_name, "chunks": chunk_count}
    except Exception as e:
        return {"success": False, "error": f"解析失败: {e}"}


@app.get("/api/knowledge/docs")
async def api_list_docs():
    return list_documents()


# ═══════════════════════════════════════════════════════════════════════════════
# 路由：统计数据 & 文件下载
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/api/stats")
async def api_stats():
    return {
        "agents_count": len([k for k in AGENT_CONFIGS if k != "skill_engineer"]),
        "skills_count": len(_skill_registry),
        "docs_count": len(_uploaded_docs),
        "workspace_files": len(list(WORKSPACE_DIR.glob("*"))) if WORKSPACE_DIR.exists() else 0,
    }


@app.get("/api/workspace/files")
async def api_list_workspace_files():
    if not WORKSPACE_DIR.exists():
        return []
    files = []
    for f in sorted(WORKSPACE_DIR.iterdir()):
        if f.is_file() and not f.name.startswith("."):
            files.append({
                "name": f.name,
                "size": f.stat().st_size,
                "modified": f.stat().st_mtime,
            })
    return files


@app.get("/api/workspace/files/{filename}")
async def api_download_workspace_file(filename: str):
    filepath = WORKSPACE_DIR / filename
    if not filepath.exists() or not filepath.is_file():
        return {"error": "文件不存在"}
    return FileResponse(filepath, filename=filename)


# ═══════════════════════════════════════════════════════════════════════════════
# 路由：Workflow 工作流
# ═══════════════════════════════════════════════════════════════════════════════

class WorkflowRunRequest(BaseModel):
    input: str

WORKFLOW_TEMPLATES = {
    "doc_pipeline": {
        "id": "doc_pipeline",
        "name": "文档处理流水线",
        "description": "上传文档 → 解析内容 → AI 分析摘要 → 生成报告",
        "steps": ["文档解析", "知识入库", "AI 分析摘要", "生成报告"],
        "icon": "FileText",
        "agents": ["a2", "a1", "a3"],
    },
    "data_pipeline": {
        "id": "data_pipeline",
        "name": "数据分析流水线",
        "description": "接收数据需求 → 数据查询 → 统计分析 → 图表生成 → 导出报告",
        "steps": ["需求理解", "数据查询", "统计分析", "生成报告"],
        "icon": "BarChart3",
        "agents": ["a1", "a3"],
    },
    "ppt_pipeline": {
        "id": "ppt_pipeline",
        "name": "PPT 生成流水线",
        "description": "描述主题 → AI 生成大纲 → 逐页生成内容 → 输出 .pptx 文件",
        "steps": ["主题分析", "大纲生成", "内容填充", "生成文件"],
        "icon": "Presentation",
        "agents": ["a7"],
    },
}


@app.get("/api/workflows")
async def api_list_workflows():
    return [
        {
            "id": wf["id"],
            "name": wf["name"],
            "description": wf["description"],
            "steps": wf["steps"],
            "icon": wf["icon"],
        }
        for wf in WORKFLOW_TEMPLATES.values()
    ]


@app.post("/api/workflows/{workflow_id}/run")
async def api_run_workflow(workflow_id: str, request: WorkflowRunRequest):
    wf = WORKFLOW_TEMPLATES.get(workflow_id)
    if not wf:
        return {"error": f"工作流 [{workflow_id}] 不存在"}

    steps = wf["steps"]
    agent_ids = wf["agents"]

    async def generate():
        user_input = request.input
        _heading_sep = '\n\n---\n\n'

        if workflow_id == "ppt_pipeline":
            yield f"data: {json.dumps({'content': '[STEP:1]', 'done': False})}\n\n"
            _msg = '正在分析主题并生成 PPT...\n\n'
            yield f"data: {json.dumps({'content': _msg, 'done': False})}\n\n"

            ppt_agent = get_agent("a7")
            if ppt_agent:
                prompt = "请根据以下需求生成一个 PPT 文件：" + user_input + "\n请使用 python-pptx 生成 .pptx 文件并保存。"
                try:
                    for step_idx in range(1, len(steps) + 1):
                        _step = f'[STEP:{step_idx}]'
                        yield f"data: {json.dumps({'content': _step, 'done': False})}\n\n"
                    async for chunk in ppt_agent.arun(prompt, stream=True):
                        if chunk.content:
                            yield f"data: {json.dumps({'content': chunk.content, 'done': False})}\n\n"
                except Exception as e:
                    _err = f'\n\n执行出错: {e}'
                    yield f"data: {json.dumps({'content': _err, 'done': False})}\n\n"

        elif workflow_id == "data_pipeline":
            data_agent = get_agent("a1")
            code_agent = get_agent("a3")

            for step_idx, step_name in enumerate(steps, 1):
                _step = f'[STEP:{step_idx}]'
                yield f"data: {json.dumps({'content': _step, 'done': False})}\n\n"
                _heading = f'### 步骤 {step_idx}: {step_name}\n\n'
                yield f"data: {json.dumps({'content': _heading, 'done': False})}\n\n"

                if step_idx <= 3 and data_agent:
                    prompt = f"[{step_name}] 用户需求：{user_input}"
                    try:
                        async for chunk in data_agent.arun(prompt, stream=True):
                            if chunk.content:
                                yield f"data: {json.dumps({'content': chunk.content, 'done': False})}\n\n"
                    except Exception as e:
                        _err = f'出错: {e}'
                        yield f"data: {json.dumps({'content': _err, 'done': False})}\n\n"
                elif step_idx == 4 and code_agent:
                    try:
                        async for chunk in code_agent.arun(
                            f"请根据数据分析结果生成一份总结报告：{user_input}", stream=True
                        ):
                            if chunk.content:
                                yield f"data: {json.dumps({'content': chunk.content, 'done': False})}\n\n"
                    except Exception as e:
                        _err = f'出错: {e}'
                        yield f"data: {json.dumps({'content': _err, 'done': False})}\n\n"

                yield f"data: {json.dumps({'content': _heading_sep, 'done': False})}\n\n"

        elif workflow_id == "doc_pipeline":
            knowledge_agent = get_agent("a2")
            data_agent = get_agent("a1")

            for step_idx, step_name in enumerate(steps, 1):
                _step = f'[STEP:{step_idx}]'
                yield f"data: {json.dumps({'content': _step, 'done': False})}\n\n"
                _heading = f'### 步骤 {step_idx}: {step_name}\n\n'
                yield f"data: {json.dumps({'content': _heading, 'done': False})}\n\n"

                agent_to_use = knowledge_agent if step_idx <= 2 else data_agent
                if agent_to_use:
                    prompt = f"[{step_name}] 用户需求：{user_input}"
                    try:
                        async for chunk in agent_to_use.arun(prompt, stream=True):
                            if chunk.content:
                                yield f"data: {json.dumps({'content': chunk.content, 'done': False})}\n\n"
                    except Exception as e:
                        _err = f'出错: {e}'
                        yield f"data: {json.dumps({'content': _err, 'done': False})}\n\n"

                yield f"data: {json.dumps({'content': _heading_sep, 'done': False})}\n\n"

        yield f"data: {json.dumps({'content': '', 'done': True})}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )
